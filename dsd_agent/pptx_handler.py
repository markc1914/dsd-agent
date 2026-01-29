"""PowerPoint handler for reading and writing DSD documents."""

from dataclasses import dataclass
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu


@dataclass
class Placeholder:
    """Represents a Lorem Ipsum placeholder in the presentation."""
    slide_index: int
    slide_title: str
    shape_name: str
    text: str
    left: float  # inches
    top: float  # inches
    width: float
    height: float
    row_group: int  # logical row grouping


@dataclass
class ArchitectureSlide:
    """Represents a slide containing architecture placeholders."""
    index: int  # 0-based
    title: str
    placeholders: list[Placeholder]


class DSDDocument:
    """Handler for Digital Solutioning Document PowerPoint files."""

    EMU_PER_INCH = 914400

    def __init__(self, pptx_path: str | Path):
        self.path = Path(pptx_path)
        self.prs = Presentation(str(self.path))
        self._architecture_slides: list[ArchitectureSlide] | None = None

    @property
    def slide_count(self) -> int:
        return len(self.prs.slides)

    def _get_slide_title(self, slide) -> str:
        """Extract the title from a slide."""
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text and len(text) < 100:
                    # Check if it looks like a title (placeholder or first text)
                    if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                        return text
        # Fallback to first short text
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text and len(text) < 80:
                    return text
        return "Untitled"

    def _emu_to_inches(self, emu: int) -> float:
        """Convert EMUs to inches."""
        return emu / self.EMU_PER_INCH

    def _assign_row_groups(self, placeholders: list[Placeholder]) -> list[Placeholder]:
        """Assign logical row groups based on vertical position."""
        if not placeholders:
            return placeholders

        # Sort by top position
        sorted_ph = sorted(placeholders, key=lambda p: p.top)

        row_group = 0
        current_top = sorted_ph[0].top

        for ph in sorted_ph:
            if abs(ph.top - current_top) > 0.5:  # More than 0.5 inch difference = new row
                row_group += 1
                current_top = ph.top
            ph.row_group = row_group

        return placeholders

    def find_architecture_slides(self) -> list[ArchitectureSlide]:
        """Find all slides containing Lorem Ipsum placeholders in architecture sections."""
        if self._architecture_slides is not None:
            return self._architecture_slides

        architecture_slides = []

        for idx, slide in enumerate(self.prs.slides):
            title = self._get_slide_title(slide)
            placeholders = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if "lorem" in text.lower():
                        ph = Placeholder(
                            slide_index=idx,
                            slide_title=title,
                            shape_name=shape.name,
                            text=text,
                            left=self._emu_to_inches(shape.left),
                            top=self._emu_to_inches(shape.top),
                            width=self._emu_to_inches(shape.width),
                            height=self._emu_to_inches(shape.height),
                            row_group=0
                        )
                        placeholders.append(ph)

            if placeholders:
                # Assign row groups and sort
                placeholders = self._assign_row_groups(placeholders)
                placeholders.sort(key=lambda p: (p.row_group, p.left))

                architecture_slides.append(ArchitectureSlide(
                    index=idx,
                    title=title,
                    placeholders=placeholders
                ))

        self._architecture_slides = architecture_slides
        return architecture_slides

    def update_placeholder(self, slide_index: int, shape_name: str, new_text: str) -> bool:
        """Update a specific placeholder with new text."""
        if slide_index < 0 or slide_index >= len(self.prs.slides):
            return False

        slide = self.prs.slides[slide_index]

        for shape in slide.shapes:
            if shape.name == shape_name and shape.has_text_frame:
                # Preserve formatting by updating paragraph text
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = new_text
                        return True
                # Fallback: set the first paragraph
                if shape.text_frame.paragraphs:
                    shape.text_frame.paragraphs[0].runs[0].text = new_text
                    return True
        return False

    def update_placeholders_batch(self, updates: dict[tuple[int, str], str]) -> int:
        """
        Update multiple placeholders at once.

        Args:
            updates: Dict mapping (slide_index, shape_name) to new_text

        Returns:
            Number of successful updates
        """
        count = 0
        for (slide_index, shape_name), new_text in updates.items():
            if self.update_placeholder(slide_index, shape_name, new_text):
                count += 1
        return count

    def save(self, output_path: str | Path | None = None) -> Path:
        """Save the presentation to a file."""
        if output_path is None:
            # Create output filename with _populated suffix
            output_path = self.path.parent / f"{self.path.stem}_populated{self.path.suffix}"
        else:
            output_path = Path(output_path)

        self.prs.save(str(output_path))
        return output_path

    def get_slide_summary(self) -> str:
        """Get a summary of architecture slides for display."""
        slides = self.find_architecture_slides()
        lines = [f"Found {len(slides)} slides with Lorem Ipsum placeholders:\n"]

        for slide in slides:
            lines.append(f"  Slide {slide.index + 1}: {slide.title}")
            lines.append(f"    - {len(slide.placeholders)} placeholders")

            # Group by row
            rows = {}
            for ph in slide.placeholders:
                if ph.row_group not in rows:
                    rows[ph.row_group] = []
                rows[ph.row_group].append(ph)

            for row_num, phs in sorted(rows.items()):
                lines.append(f"    - Row {row_num + 1}: {len(phs)} boxes")

        return "\n".join(lines)
